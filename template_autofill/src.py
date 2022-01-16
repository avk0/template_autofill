import copy
import os
import shutil

import pandas as pd
from pptx import Presentation


def read_presentation(path):
    return Presentation(path)


def save_presentation(prs, path):
    prs.save(path)


def read_data(path):
    df = pd.read_excel(path)
    return df
    
    
def replace_text_retaining_initial_formatting(shape, new_text):
    paragraph = shape.text_frame.paragraphs[0]
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text


    
def duplicate_slide(pres, index):
    template = pres.slides[index]
    blank_slide_layout = pres.slide_layouts[0]
    copied_slide = pres.slides.add_slide(blank_slide_layout)
    
    # remove empty shapes from slide master
    for shp in copied_slide.shapes:
        if shp.text == '':
            sp = shp._sp
            sp.getparent().remove(sp)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return copied_slide


def fill_pres_with_data(pres, data):
    prs = copy.deepcopy(pres)
    for index, row in data.iterrows():
        slide = duplicate_slide(prs, 0)
        for shape in slide.shapes:
            try:
                replace_text_retaining_initial_formatting(shape, row[shape.text])
            except Exception as e:
                pass
                #print(e)
    return prs
    

def fill_sep_pres_with_data(pres, data, path):
    for index, row in data.iterrows():
        prs = copy.deepcopy(pres)
        slide = prs.slides[0]
        #slide = duplicate_slide(prs, 0)
        for shape in slide.shapes:
            try:
                replace_text_retaining_initial_formatting(shape, row[shape.text])
            except Exception as e:
                pass
                #print(e)
        save_presentation(prs, os.path.join(path, f'{index}.pptx'))
    return pres
    
    
def save_files_as_zip(dir_name, zip_name):
    shutil.make_archive(zip_name, 'zip', dir_name)